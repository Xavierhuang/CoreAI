#!/usr/bin/env python3
"""Build the 6-minute lightning-talk deck:
    'How Much Can On-Device AI Models Do on a Mac?'

6 slides. ~60 seconds each. Speaker notes include timing.

Run with:
    /tmp/pptx-venv/bin/python build_core_ai_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- Palette ----------
BG       = RGBColor(0x0B, 0x0F, 0x14)
PANEL    = RGBColor(0x14, 0x1A, 0x22)
PANEL_2  = RGBColor(0x1A, 0x22, 0x2C)
INK      = RGBColor(0xEC, 0xEF, 0xF4)
INK_DIM  = RGBColor(0x9A, 0xA4, 0xB2)
ACCENT   = RGBColor(0x0A, 0x84, 0xFF)
GREEN    = RGBColor(0x32, 0xD7, 0x4D)
YELLOW   = RGBColor(0xFF, 0xD6, 0x0A)
ORANGE   = RGBColor(0xFF, 0x9F, 0x0A)
RED      = RGBColor(0xFF, 0x45, 0x3A)

OUTPUT = "/Users/weijiahuang/Desktop/ide/On-Device-AI-Mac-Talk.pptx"

# ---------- Helpers ----------

def new_deck():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Inches(13.333), Inches(7.5))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg

def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font="Helvetica Neue"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_bullets(slide, items, left, top, width, height, *,
                size=22, color=INK, font="Helvetica Neue",
                spacing=Pt(14)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = spacing
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def add_accent_bar(slide, color=ACCENT, top=Inches(1.4), height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), top, Inches(0.6), height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar

def add_title(slide, title, color=INK, accent=ACCENT):
    add_text(slide, title, Inches(0.6), Inches(0.55),
             Inches(12), Inches(0.8), size=36, bold=True, color=color)
    add_accent_bar(slide, color=accent)

def add_footer(slide, idx, total,
               footer_text="How Much Can On-Device AI Do on a Mac?"):
    add_text(slide, footer_text,
             Inches(0.6), Inches(7.1), Inches(11), Inches(0.3),
             size=10, color=INK_DIM)
    add_text(slide, f"{idx} / {total}",
             Inches(11.7), Inches(7.1), Inches(1.0), Inches(0.3),
             size=10, color=INK_DIM, align=PP_ALIGN.RIGHT)

def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    add_bg(s)
    return s

def table(slide, rows, left, top, width, height, *,
          header_fill=ACCENT, header_ink=INK,
          row_fill=PANEL, row_ink=INK,
          stripe_fill=None, header_size=14, body_size=13,
          col_color_map=None):
    n_rows = len(rows)
    n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            tint = None
            if col_color_map and r > 0 and c in col_color_map:
                for sub, color in col_color_map[c].items():
                    if sub in cell_text:
                        tint = color
                        break
            if r == 0:
                cell.fill.fore_color.rgb = header_fill
                ink = header_ink
                bold = True
                size = header_size
            else:
                if tint is not None:
                    cell.fill.fore_color.rgb = tint
                    ink = INK
                elif stripe_fill is not None and (r % 2 == 0):
                    cell.fill.fore_color.rgb = stripe_fill
                    ink = row_ink
                else:
                    cell.fill.fore_color.rgb = row_fill
                    ink = row_ink
                bold = False
                size = body_size
            tf = cell.text_frame
            tf.clear()
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Helvetica Neue"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = ink
    return tbl_shape

# ---------- Slides ----------

prs = new_deck()
TOTAL = 6

# 1. Title
s = blank_slide(prs)
add_text(s, "How Much Can",
         Inches(0.6), Inches(1.5), Inches(12), Inches(1.0),
         size=64, bold=True, color=INK)
add_text(s, "On-Device AI",
         Inches(0.6), Inches(2.4), Inches(12), Inches(1.0),
         size=64, bold=True, color=ACCENT)
add_text(s, "Do on a Mac?",
         Inches(0.6), Inches(3.3), Inches(12), Inches(1.0),
         size=64, bold=True, color=INK)
add_text(s,
         "Hands-on data from a 16 GB MacBook · 5 backends compared",
         Inches(0.6), Inches(4.8), Inches(12), Inches(0.6),
         size=20, color=INK_DIM)
add_text(s, "6-minute lightning talk · 2026",
         Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         size=12, color=INK_DIM)
add_notes(s,
"[~20 seconds] Open with the question every dev is quietly asking: can I "
"actually replace ChatGPT with something on my laptop? Tonight I'm "
"answering with real numbers from a 16 GB MacBook — the median dev "
"machine, not some 96 GB monster. Five backends, same prompts.")

# 2. My setup
s = blank_slide(prs)
add_title(s, "My setup")
add_bullets(s, [
    "16 GB MacBook · macOS 26.5 (Tahoe-class)",
    "Built a SwiftUI IDE with 5 swappable AI backends in one picker",
    "Claude (cloud) · Apple Intelligence · Apple PCC · Ollama · Core AI",
    "Same prompts, same project folder — apples-to-apples",
], Inches(0.7), Inches(2.0), Inches(12), Inches(3.0), size=22)
add_text(s,
         "→ Lets me ask the same question through every available stack and time it.",
         Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         size=18, color=ACCENT, bold=True)
add_footer(s, 2, TOTAL)
add_notes(s,
"[~45 seconds] Quick credibility. I built a SwiftUI IDE that has all five "
"backends side-by-side. Optional: flash the IDE for 2-3 seconds, show the "
"picker. Don't dwell — the data slide is what matters. Move on.")

# 3. Speed reality (THE DATA SLIDE)
s = blank_slide(prs)
add_title(s, "Speed reality — same prompt, 5 backends", accent=GREEN)
rows = [
    ["Backend",                             "First-token", "Total reply",  "Quality",   "Runs on"],
    ["Claude Sonnet (cloud)",               "~600 ms",     "~3 s",         "Excellent", "Anthropic"],
    ["Apple Intelligence (~3B on-device)",  "~400 ms",     "~3 s",         "Decent",    "ANE + Metal"],
    ["qwen2.5-coder:1.5b (Ollama)",         "~1.5 s",      "~5–8 s",       "Good",      "Metal"],
    ["qwen3:0.6b (Ollama)",                 "~800 ms",     "~3 s",         "Mediocre",  "Metal"],
    ["deepseek-coder:6.7b (Ollama, 16 GB)", "~30 s",       "~3 MINUTES",   "Unusable",  "CPU"],
]
table(s, rows, Inches(0.3), Inches(1.9), Inches(12.7), Inches(3.8),
      header_fill=GREEN, header_ink=BG,
      row_fill=PANEL, stripe_fill=PANEL_2,
      header_size=15, body_size=14,
      col_color_map={
          4: {"ANE": ACCENT, "Metal": GREEN, "CPU": RED, "Anthropic": INK_DIM},
          3: {"Excellent": GREEN, "Good": GREEN, "Decent": YELLOW, "Mediocre": YELLOW, "Unusable": RED},
      })
add_text(s,
         "Small models (≤3B) on Apple Silicon are basically as fast as cloud calls.",
         Inches(0.3), Inches(6.0), Inches(12.7), Inches(0.5),
         size=17, color=GREEN, bold=True)
add_text(s,
         "The 7B model on 16 GB ran for 3 minutes. That's the wall.",
         Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.5),
         size=17, color=RED, bold=True)
add_footer(s, 3, TOTAL)
add_notes(s,
"[~90 seconds] THE money slide. Walk left to right. Cloud Claude is the "
"baseline — ~3 seconds for a real answer. Apple's on-device, also 3 "
"seconds. qwen 1.5b coder via Ollama, 5-8 seconds, and it actually uses "
"tools. Tiny qwen 0.6b, fast but dumb. "
"Then point at the bottom row: deepseek 6.7B on a 16 GB Mac took 3 MINUTES "
"for the same 'hello' prompt. Same hardware, same Ollama, same model "
"format. Why? Next slide.")

# 4. RAM is the wall
s = blank_slide(prs)
add_title(s, "Why the cliff? RAM is the wall.", accent=ORANGE)
add_bullets(s, [
    "Apple Silicon shares ONE memory pool — CPU, GPU, and Neural Engine",
    "16 GB total · Chrome eats ~6 GB · Slack ~2 · IDE ~3 · Discord ~1",
    "When free RAM drops below the model's size, Ollama silently demotes it to CPU",
    "Same model: 5 seconds on Metal → 3 minutes on CPU. Same hardware.",
], Inches(0.7), Inches(2.0), Inches(12), Inches(3.5), size=21)
add_text(s,
         "Rule of thumb on 16 GB:  ≤3B = comfortable · 7B = painful · 13B+ = nope.",
         Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
         size=18, color=ORANGE, bold=True)
add_text(s,
         "Always quit Chrome before you benchmark.",
         Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
         size=14, color=INK_DIM)
add_footer(s, 4, TOTAL)
add_notes(s,
"[~60 seconds] The thing no review tells you. Apple Silicon's unified "
"memory is great until you run out of it. When the model can't fit "
"alongside everything else you have open, Ollama silently puts it on the "
"CPU and you go from 'instant' to 'walk away and come back'. "
"The rule: on a 16 GB Mac, stay at 3B parameters or below. Above that, "
"you need 24 GB to be safe.")

# 5. The verdict
s = blank_slide(prs)
add_title(s, "The 2026 verdict on a 16 GB Mac", accent=ACCENT)
verdicts = [
    ("Replace ChatGPT entirely for general chat?",         "No, not yet",            RED),
    ("Code completion, specialized small models?",          "Yes — better than cloud", GREEN),
    ("Privacy-critical assistant (legal, medical, notes)?", "Yes — only real option",  GREEN),
    ("Offline / airplane / no-network use cases?",          "Yes — finally usable",    GREEN),
    ("Long documents, complex multi-step reasoning?",       "Cloud or PCC for now",   YELLOW),
]
y = Inches(1.95)
for q, verdict, color in verdicts:
    add_text(s, q,
             Inches(0.7), y, Inches(8.3), Inches(0.6),
             size=20, color=INK)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(9.2), y + Inches(0.04),
                              Inches(3.7), Inches(0.55))
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    tf = pill.text_frame
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = verdict
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = BG if color in (GREEN, YELLOW) else INK
    y += Inches(0.85)
add_footer(s, 5, TOTAL)
add_notes(s,
"[~90 seconds] The honest verdict. On 16 GB, on-device can't replace "
"ChatGPT for general chat — yet. But for specialized tasks, privacy, and "
"offline use, it's not just viable, it's BETTER than cloud. "
"For long-context or hard reasoning, you still want cloud or Apple's "
"PCC. The pattern in 2026 is hybrid — small fast on-device by default, "
"cloud fallback for the hard stuff.")

# 6. Summary + Q&A
s = blank_slide(prs)
add_text(s, "The takeaway",
         Inches(0.6), Inches(0.55), Inches(12), Inches(0.8),
         size=36, bold=True, color=INK)
add_accent_bar(s)
add_bullets(s, [
    "On a 16 GB Mac, ≤3B models are FAST and usable for the right tasks",
    "RAM is the wall, not the model — bigger models silently fall back to CPU",
    "Hybrid (small on-device + cloud / PCC fallback) is the realistic architecture",
    "Don't trust demos on 96 GB Macs — test on the machine your users own",
], Inches(0.7), Inches(1.8), Inches(12), Inches(3.5), size=22)
add_text(s, "Questions?",
         Inches(0.6), Inches(5.7), Inches(12), Inches(0.9),
         size=52, bold=True, color=ACCENT)
add_text(s, "Demo IDE: github.com/<you>/ide   (placeholder)",
         Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         size=14, color=INK_DIM)
add_footer(s, 6, TOTAL)
add_notes(s,
"[~30 seconds] Recap and stop. Four bullets: small models are fast on "
"Apple Silicon, RAM is the binding constraint, hybrid is the real answer, "
"test on median hardware. Open for questions. Total: ~5:35 leaving ~25s "
"buffer.")

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}  (target: 6 for a 6-minute talk)")
