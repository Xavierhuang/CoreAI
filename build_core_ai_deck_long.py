#!/usr/bin/env python3
"""Build the LONG (16-slide) version of the on-device-AI talk.

Use this for a 20-30 minute slot. For the 6-minute lightning version, see
build_core_ai_deck.py.

Run with:
    /tmp/pptx-venv/bin/python build_core_ai_deck_long.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- Palette ----------
BG       = RGBColor(0x0B, 0x0F, 0x14)
PANEL    = RGBColor(0x14, 0x1A, 0x22)
PANEL_2  = RGBColor(0x1A, 0x22, 0x2C)
INK      = RGBColor(0xEC, 0xEF, 0xF4)
INK_DIM  = RGBColor(0x9A, 0xA4, 0xB2)
ACCENT   = RGBColor(0x0A, 0x84, 0xFF)
GREEN    = RGBColor(0x32, 0xD7, 0x4D)
YELLOW   = RGBColor(0xFF, 0xD6, 0x0A)
ORANGE   = RGBColor(0xFF, 0x9F, 0x0A)
RED      = RGBColor(0xFF, 0x45, 0x3A)
PURPLE   = RGBColor(0xBF, 0x5A, 0xF2)
TEAL     = RGBColor(0x40, 0xC8, 0xE0)

OUTPUT = "/Users/weijiahuang/Desktop/ide/On-Device-AI-Mac-Talk-Long.pptx"

# ---------- Helpers ----------

def new_deck():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Inches(13.333), Inches(7.5))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg

def add_text(slide, text, left, top, width, height, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font="Helvetica Neue"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_bullets(slide, items, left, top, width, height, *,
                size=18, color=INK, font="Helvetica Neue",
                spacing=Pt(8)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = spacing
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb

def add_accent_bar(slide, color=ACCENT, top=Inches(1.4), height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), top, Inches(0.6), height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar

def add_title(slide, title, color=INK, accent=ACCENT):
    add_text(slide, title, Inches(0.6), Inches(0.55),
             Inches(12), Inches(0.8), size=32, bold=True, color=color)
    add_accent_bar(slide, color=accent)

def add_footer(slide, idx, total,
               footer_text="How Much Can On-Device AI Models Do on a Mac"):
    add_text(slide, footer_text,
             Inches(0.6), Inches(7.1), Inches(11), Inches(0.3),
             size=10, color=INK_DIM)
    add_text(slide, f"{idx} / {total}",
             Inches(11.7), Inches(7.1), Inches(1.0), Inches(0.3),
             size=10, color=INK_DIM, align=PP_ALIGN.RIGHT)

def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    add_bg(s)
    return s

def table(slide, rows, left, top, width, height, *,
          header_fill=ACCENT, header_ink=INK,
          row_fill=PANEL, row_ink=INK,
          stripe_fill=None, header_size=14, body_size=12,
          col_color_map=None):
    n_rows = len(rows)
    n_cols = len(rows[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            tint = None
            if col_color_map and r > 0 and c in col_color_map:
                for sub, color in col_color_map[c].items():
                    if sub in cell_text:
                        tint = color
                        break
            if r == 0:
                cell.fill.fore_color.rgb = header_fill
                ink = header_ink
                bold = True
                size = header_size
            else:
                if tint is not None:
                    cell.fill.fore_color.rgb = tint
                    ink = INK
                elif stripe_fill is not None and (r % 2 == 0):
                    cell.fill.fore_color.rgb = stripe_fill
                    ink = row_ink
                else:
                    cell.fill.fore_color.rgb = row_fill
                    ink = row_ink
                bold = False
                size = body_size
            tf = cell.text_frame
            tf.clear()
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Helvetica Neue"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = ink
    return tbl_shape

# ---------- Slides ----------

prs = new_deck()
TOTAL = 16

# 1. Title
s = blank_slide(prs)
add_text(s, "How Much Can On-Device",
         Inches(0.6), Inches(1.8), Inches(12), Inches(1.0),
         size=58, bold=True, color=INK)
add_text(s, "AI Models Do on a Mac?",
         Inches(0.6), Inches(2.7), Inches(12), Inches(1.0),
         size=58, bold=True, color=INK)
add_text(s,
         "A hands-on look at what's actually possible in 2026 — and what isn't.",
         Inches(0.6), Inches(4.2), Inches(12), Inches(0.8),
         size=22, color=INK_DIM)
add_text(s, "Tested live on a 16 GB MacBook · 5 backends · same prompts",
         Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         size=14, color=ACCENT)
add_notes(s,
"Hi everyone. The question I want to answer is the one every developer is "
"quietly asking: can I actually run useful AI on my Mac, instead of paying "
"OpenAI every time? I'll show what works, what's painful, what's not "
"happening — all from real testing on a 16 GB MacBook.")

# 2. The question
s = blank_slide(prs)
add_title(s, "The question every dev is asking")
add_bullets(s, [
    "Can I replace ChatGPT / Claude with something running on MY laptop?",
    "Without an API key, without a subscription, without a network?",
    "How close are we to that today — and on what hardware?",
    "What's the honest gap between 'demo' and 'I'd ship this in production'?",
], Inches(0.7), Inches(2.0), Inches(12), Inches(3.5), size=22)
add_text(s, "Today I'll answer with real numbers, not vibes.",
         Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
         size=18, color=ACCENT, bold=True)
add_footer(s, 2, TOTAL)
add_notes(s,
"Most posts about on-device AI either say 'it's amazing' or 'it's "
"useless'. Both are wrong. The truth depends on what task, what hardware, "
"what model. Decision framework with actual data follows.")

# 3. Your options on a Mac
s = blank_slide(prs)
add_title(s, "Your options on a Mac in 2026")
rows = [
    ["Where it runs",    "What it is",                                                   "Network?"],
    ["FoundationModels", "Apple's pre-trained on-device LLM (~3B). Ships in macOS 26+.", "No"],
    ["Apple PCC",        "Apple's bigger model in Private Cloud Compute. macOS 27+.",     "Yes (attested)"],
    ["Core AI",          ".aimodel files loaded by Apple's runtime. macOS 27+.",          "No"],
    ["Ollama",           "llama.cpp daemon + GGUF models you `ollama pull`.",             "No"],
    ["MLX",              "Apple-Silicon-tuned ML library for custom inference.",          "No"],
    ["Cloud APIs",       "Claude / OpenAI / Gemini. The 'control' in our comparison.",   "Yes"],
]
table(s, rows, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.5),
      header_fill=ACCENT, row_fill=PANEL, stripe_fill=PANEL_2,
      header_size=15, body_size=13)
add_text(s, "I built one app with 5 of these as switchable backends, to compare apples-to-apples.",
         Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
         size=14, color=INK_DIM)
add_footer(s, 3, TOTAL)
add_notes(s,
"Most devs don't realize all of these exist. I built a side-by-side "
"comparison in a SwiftUI IDE with 5 of these behind a picker.")

# 4. What "on-device" means
s = blank_slide(prs)
add_title(s, "What 'on-device' actually means in 2026")
add_bullets(s, [
    "Apple Silicon: unified memory shared between CPU, GPU (Metal), and ANE (Neural Engine)",
    "Models live in that pool — RAM is the binding constraint, not GPU VRAM",
    "8 GB / 16 GB / 24 GB / 32 GB / 48 GB / 64 GB / 128 GB across the M-series lineup",
    "Models range from 0.5B (fits in 1 GB) to 70B+ (needs 32-48 GB minimum)",
    "Real perf depends on: model size · quantization · ANE vs GPU vs CPU dispatch",
], Inches(0.7), Inches(1.9), Inches(12), Inches(4.5), size=18)
add_footer(s, 4, TOTAL)
add_notes(s,
"Apple Silicon shares one memory pool across CPU/GPU/ANE. So your laptop's "
"total RAM is the hard ceiling for models you can load.")

# 5. RAM tier table
s = blank_slide(prs)
add_title(s, "The hardware wall: RAM tier = which models you can run", accent=ORANGE)
rows = [
    ["Mac RAM",    "Comfortable",                          "Painful",                     "Won't fit"],
    ["8 GB",       "0.5–1B (chat toys)",                   "3B",                          "7B+"],
    ["16 GB",      "0.5–3B (Apple Intelligence sized)",    "7B (drops to CPU)",           "13B+"],
    ["24 GB",      "0.5–7B (most OSS models)",             "13B",                          "30B+"],
    ["32 GB",      "0.5–13B (Llama 3 / Qwen 7B/14B)",      "32B",                          "70B"],
    ["48 GB",      "up to 30B comfortably",                 "70B Q4",                      "—"],
    ["64–128 GB",  "70B comfortably; multiple at once",    "—",                            "—"],
]
table(s, rows, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.4),
      header_fill=ORANGE, row_fill=PANEL, stripe_fill=PANEL_2,
      header_size=14, body_size=13,
      col_color_map={
          1: {"comfortably": GREEN},
      })
add_text(s,
         "→ 16 GB is the median dev machine. It's also where the gap with cloud is widest.",
         Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
         size=15, color=ORANGE, bold=True)
add_footer(s, 5, TOTAL)
add_notes(s,
"Most reviews test on 96 GB M3 Max. Reality is 16 GB is the median. At "
"16 GB you can run small models well, 7B with pain, 13B+ is impossible.")

# 6. My setup
s = blank_slide(prs)
add_title(s, "My setup for this talk")
add_bullets(s, [
    "MacBook with 16 GB unified memory (Apple Silicon)",
    "macOS 26.5 (Tahoe is 27 — Apple Intelligence works, PCC + Core AI runtime don't yet)",
    "Custom SwiftUI IDE with a 5-backend chat picker: Claude · Apple · PCC · Ollama · Core AI",
    "Models tested: Apple's built-in (~3B), qwen2.5:0.5b, qwen3:0.6b, qwen2.5-coder:1.5b, deepseek-coder:6.7b",
    "All asking the same questions over the same project folder",
], Inches(0.7), Inches(1.9), Inches(12), Inches(4.5), size=18)
add_text(s,
         "Goal: not benchmarks. 'Does this answer feel useful, and how long did I wait?'",
         Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
         size=15, color=INK_DIM)
add_footer(s, 6, TOTAL)
add_notes(s,
"DEMO opportunity. Show the IDE, the picker, the same prompt going through "
"Claude vs Apple vs Ollama.")

# 7. Speed reality check
s = blank_slide(prs)
add_title(s, "Speed reality check — same prompt, different backends", accent=GREEN)
rows = [
    ["Backend / Model",                  "First-token",  "Total reply (~150 tok)", "Quality",          "Where it runs"],
    ["Claude Sonnet (cloud)",            "~600 ms",      "~3 s",                    "Excellent",        "Anthropic API"],
    ["Apple Intelligence on-device (3B)","~400 ms",      "~3 s",                    "Decent — chatty",  "ANE + Metal"],
    ["qwen2.5:0.5b (Ollama)",            "~500 ms",      "~2 s",                    "Mediocre",         "Metal"],
    ["qwen3:0.6b (Ollama)",              "~800 ms",      "~3 s",                    "Mediocre+",        "Metal"],
    ["qwen2.5-coder:1.5b (Ollama)",      "~1.5 s",       "~5–8 s",                  "Good, tool-using", "Metal"],
    ["deepseek-coder:6.7b (Ollama, 16GB)","~30 s",        "~3 min",                  "Good but unusable","CPU (RAM-evicted)"],
]
table(s, rows, Inches(0.3), Inches(1.85), Inches(12.7), Inches(4.6),
      header_fill=GREEN, header_ink=BG,
      row_fill=PANEL, stripe_fill=PANEL_2,
      header_size=13, body_size=12,
      col_color_map={
          4: {"ANE": ACCENT, "Metal": GREEN, "CPU": RED, "API": INK_DIM},
          3: {"Excellent": GREEN, "Good": GREEN, "Decent": YELLOW, "Mediocre": YELLOW, "unusable": RED},
      })
add_text(s,
         "Key insight: 0.5–1.5B models on Apple Silicon are FAST. 6.7B on 16 GB is unusable.",
         Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.5),
         size=15, color=GREEN, bold=True)
add_footer(s, 7, TOTAL)
add_notes(s,
"Cliff at the bottom: deepseek 6.7B on 16 GB takes literally 3 minutes "
"because it's running on CPU. Same model on a 32 GB Mac would be 5-10 sec.")

# 8. What on-device CAN do
s = blank_slide(prs)
add_title(s, "What on-device DOES well today", accent=GREEN)
add_bullets(s, [
    "Code completion / inline suggestions (specialized 1–3B coder models)",
    "Privacy-sensitive Q&A (legal, medical, personal notes that can't leave the device)",
    "Offline assistant — flights, basement, sketchy hotel wifi, oil rigs",
    "Speech transcription (Whisper variants) — same quality as cloud, no upload",
    "Image generation (Stable Diffusion via Core AI) — local Midjourney",
    "Simple agents with pre-injected context (file Q&A, journal summarization)",
    "Anything you would have called 'good' from GPT-3.5 in 2023",
], Inches(0.7), Inches(1.85), Inches(12), Inches(5.0), size=18, color=INK)
add_footer(s, 8, TOTAL)
add_notes(s,
"Code completion is the sweet spot — small specialized models, low-latency, "
"no per-token cost. Privacy is huge. Offline unlocks new app categories.")

# 9. What on-device STRUGGLES with
s = blank_slide(prs)
add_title(s, "What on-device still struggles with", accent=RED)
add_bullets(s, [
    "Complex multi-step reasoning — small models hallucinate plausible-but-wrong logic",
    "Long context (>4–16K tokens) — most local models have small windows",
    "Tool calling on tiny models — schema-aware but skips tools half the time",
    "Anything needing world knowledge breadth — small models forget the tail",
    "Frontier coding (large refactors, full-repo edits) — still cloud-only territory",
    "Anything where 'feels like ChatGPT' is the bar — small model = smaller feel",
], Inches(0.7), Inches(1.85), Inches(12), Inches(5.0), size=18, color=INK)
add_footer(s, 9, TOTAL)
add_notes(s,
"Tonight: Apple Intelligence confidently made up filenames that don't "
"exist when asked about my open project. Model isn't lying — too small to "
"ground itself in context.")

# 10. RAM pressure
s = blank_slide(prs)
add_title(s, "The hidden cost no one warns you about: RAM pressure", accent=ORANGE)
add_bullets(s, [
    "Apple Silicon shares ONE memory pool across CPU + GPU + ANE",
    "A 4 GB model loaded on Metal needs ~6 GB of free RAM (model + KV cache + buffers)",
    "Chrome with 30 tabs eats ~6 GB. Slack ~2 GB. Discord ~1 GB. Your IDE ~3 GB.",
    "Result: Ollama silently evicts the model to CPU when you're low on free RAM",
    "Same model: 5 seconds → 3 minutes. Same hardware. Same prompt.",
    "Tonight's screenshot: deepseek-coder:6.7b ran for 211 s with no answer",
], Inches(0.7), Inches(1.85), Inches(12), Inches(5.0), size=18, color=INK)
add_text(s, "Always quit Chrome before you benchmark.",
         Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
         size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_footer(s, 10, TOTAL)
add_notes(s,
"Every blog post leaves this out. Reviews benchmark on fresh boot with "
"nothing else running. Reality: your machine is already chewing 12 GB "
"before you open Ollama. Model gets demoted to CPU; interactive becomes "
"walk-away.")

# 11. Apple's stack
s = blank_slide(prs)
add_title(s, "What Apple uniquely enables (the 26/27 era)")
add_bullets(s, [
    "FoundationModels — 1 line of Swift, free, on-device, with tool calling",
    "Private Cloud Compute — when on-device too small, Apple's bigger model with attested privacy",
    "Core AI runtime — ship your own .aimodel inside your .app, no daemon, no install",
    "Hardware-software co-design: every M-series gen jumps ANE TFLOPS substantially",
    "Image inputs + structured output + tool API all in the macOS 27 SDK",
    "Hybrid pattern: small fast on-device by default, PCC fallback for hard prompts",
], Inches(0.7), Inches(1.85), Inches(12), Inches(5.0), size=18)
add_footer(s, 11, TOTAL)
add_notes(s,
"Apple's contribution: integrated stack. FoundationModels for easy "
"integration. PCC for when on-device isn't enough, with same privacy "
"guarantees.")

# 12. Open ecosystem
s = blank_slide(prs)
add_title(s, "What the open ecosystem adds")
rows = [
    ["Tool",          "What it's for",                                              "Friction"],
    ["Ollama",        "Try any GGUF model in 30 seconds. Best for experimentation.","Daemon install"],
    ["MLX",           "Apple-Silicon-tuned PyTorch-ish. Research / training.",       "Python required"],
    ["llama.cpp",     "Low-level. Cross-platform. The substrate Ollama runs on.",    "Build-it-yourself"],
    ["Hugging Face",  "The model catalog. Everything starts here.",                  "Browser tabs"],
    ["LM Studio",     "GUI for trying models. Good for non-CLI users.",              "Another app"],
]
table(s, rows, Inches(0.5), Inches(1.9), Inches(12.3), Inches(3.4),
      header_fill=PURPLE, row_fill=PANEL, stripe_fill=PANEL_2,
      header_size=14, body_size=13)
add_text(s,
         "You can mix them. My IDE ships with Apple + Ollama + Core AI side by side.",
         Inches(0.5), Inches(5.7), Inches(12), Inches(0.5),
         size=15, color=INK_DIM)
add_text(s,
         "These aren't either/or. Your app can support both Apple's stack AND open-source.",
         Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
         size=15, color=INK_DIM)
add_footer(s, 12, TOTAL)
add_notes(s,
"Right architecture for most consumer apps: FoundationModels as default, "
"Ollama or Core AI as opt-in power-user backend, cloud fallback for "
"hardest tasks.")

# 13. Verdict pills
s = blank_slide(prs)
add_title(s, "The 2026 verdict — can on-device replace cloud?", accent=TEAL)
verdicts = [
    ("Replace ChatGPT entirely on a 16 GB Mac?",        "No",                       RED),
    ("Replace it on a 64 GB+ Mac?",                      "Almost",                   YELLOW),
    ("Specialized tasks (code complete, transcribe)?",   "Yes — better than cloud",  GREEN),
    ("Privacy-critical assistant?",                      "Yes — only option",        GREEN),
    ("'Feels like ChatGPT 4o' for general chat?",        "Not yet — wait 12-18 mo",  YELLOW),
    ("Offline / airplane / no-network use cases?",       "Yes — finally",            GREEN),
]
y = Inches(2.0)
for q, verdict, color in verdicts:
    add_text(s, q,
             Inches(0.7), y, Inches(8.5), Inches(0.55),
             size=18, color=INK)
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(9.4), y + Inches(0.03),
                              Inches(3.4), Inches(0.5))
    pill.adjustments[0] = 0.5
    pill.line.fill.background()
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    tf = pill.text_frame
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = verdict
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = BG if color in (GREEN, YELLOW) else INK
    y += Inches(0.7)
add_footer(s, 13, TOTAL)
add_notes(s,
"On 16 GB: not for general chat, yes for specialized tasks. On 64 GB+: "
"almost yes. Next 12-18 mo closes most of this gap.")

# 14. Practical advice
s = blank_slide(prs)
add_title(s, "Practical advice for your next AI feature")
rows = [
    ["If your feature is…",                          "Use…",                                 "Why"],
    ["Inline code completion",                       "Small coder model on-device (Core AI / Ollama)", "Latency matters; specialized"],
    ["Spam / classification",                        "FoundationModels (free)",              "Trivial for a 3B model"],
    ["Summarize a 50-page PDF",                      "Cloud or PCC",                         "Long context = need big model"],
    ["Anything privacy-sensitive (legal/medical)",   "On-device only",                       "Liability + trust"],
    ["'Feels like ChatGPT for normal users'",        "Cloud (still)",                        "Quality gap is real"],
    ["Speech transcription",                         "Whisper via Core AI",                  "Same quality as cloud, free"],
    ["Image generation in app",                      "SD via Core AI",                       "Bundled, no daemon"],
    ["You don't know yet",                           "FoundationModels first, decide later", "Cheapest integration"],
]
table(s, rows, Inches(0.3), Inches(1.85), Inches(12.7), Inches(4.8),
      header_fill=ACCENT, row_fill=PANEL, stripe_fill=PANEL_2,
      header_size=14, body_size=12)
add_footer(s, 14, TOTAL)
add_notes(s,
"Match task to backend. Most apps end up with 2-3 backends. "
"FoundationModels is the cheapest first integration.")

# 15. What's coming
s = blank_slide(prs)
add_title(s, "What's coming — 12 to 24 months out")
add_bullets(s, [
    "M5 / M6 — bigger ANE, more memory bandwidth, more unified RAM per dollar",
    "8B-class models matching GPT-4o-level quality on many tasks (already close)",
    "Hybrid (on-device + PCC) becoming the default architecture for Apple apps",
    "Vision-language models (multimodal) on-device — image understanding, not just generation",
    "Longer context windows (32K → 128K) shrinking the 'must use cloud' set",
    "Standard ANE-tuned model bundles → less Python export friction",
], Inches(0.7), Inches(1.9), Inches(12), Inches(5.0), size=18)
add_footer(s, 15, TOTAL)
add_notes(s,
"Moving fast. Today's 'painful' on 16 GB will be smooth in 18 months — "
"hardware (M5/M6, more RAM-per-dollar) + model efficiency.")

# 16. Summary
s = blank_slide(prs)
add_text(s, "Summary",
         Inches(0.6), Inches(0.55), Inches(12), Inches(0.8),
         size=32, bold=True, color=INK)
add_accent_bar(s)
add_bullets(s, [
    "On-device on a Mac works — for the right tasks, on the right hardware",
    "RAM is the wall: 16 GB caps you at 3B comfortably, 7B with pain, 13B+ no",
    "Apple's stack (FoundationModels + PCC + Core AI) is the easiest integration",
    "Open-source (Ollama / MLX / llama.cpp) gives you variety and offline freedom",
    "The honest pattern in 2026: hybrid. Small on-device default, cloud or PCC fallback.",
    "Don't believe demos on 96 GB M3 Maxes. Test on the median machine.",
], Inches(0.7), Inches(1.7), Inches(12), Inches(3.5), size=18)
add_text(s, "Questions?",
         Inches(0.6), Inches(5.6), Inches(12), Inches(0.9),
         size=48, bold=True, color=ACCENT)
add_text(s, "Demo IDE source: github.com/<you>/ide   (placeholder)",
         Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         size=14, color=INK_DIM)
add_footer(s, 16, TOTAL)
add_notes(s,
"Recap. Hybrid is the realistic 2026 architecture. Small on-device + "
"cloud/PCC for hard prompts. Don't trust 96 GB Mac demos.")

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
